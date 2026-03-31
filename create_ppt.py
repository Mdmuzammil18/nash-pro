"""
Generate a PowerPoint presentation for the
Turbofan Engine Predictive Maintenance project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1e, 0x3a, 0x5f)
BLUE      = RGBColor(0x2d, 0x5a, 0x9e)
LIGHT_BLU = RGBColor(0xdb, 0xe4, 0xfe)
GREEN     = RGBColor(0x16, 0xa3, 0x4a)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY      = RGBColor(0xf5, 0xf7, 0xfa)
DARK_GRAY = RGBColor(0x4b, 0x55, 0x63)
RED       = RGBColor(0xc0, 0x39, 0x2b)
ORANGE    = RGBColor(0xe6, 0x7e, 0x22)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def slide_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def divider(slide, y, color=BLUE):
    add_rect(slide, Inches(0.5), y, Inches(12.33), Inches(0.04), color)

def bullet_box(slide, items, x, y, w, h, font_size=16, color=WHITE, spacing=Pt(6)):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txb

def add_card(slide, x, y, w, h, title, value, title_color=LIGHT_BLU, value_color=WHITE):
    add_rect(slide, x, y, w, h, BLUE)
    add_text_box(slide, title, x + Inches(0.15), y + Inches(0.1),
                 w - Inches(0.3), Inches(0.35), font_size=11, color=title_color)
    add_text_box(slide, value, x + Inches(0.15), y + Inches(0.4),
                 w - Inches(0.3), Inches(0.5), font_size=22, bold=True, color=value_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)

# accent bar left
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

# big title
add_text_box(sl, "✈️  Turbofan Engine", Inches(0.5), Inches(1.6), Inches(9), Inches(1.1),
             font_size=48, bold=True, color=WHITE)
add_text_box(sl, "Predictive Maintenance", Inches(0.5), Inches(2.6), Inches(9), Inches(1.1),
             font_size=48, bold=True, color=LIGHT_BLU)

divider(sl, Inches(3.85), GREEN)

add_text_box(sl, "NASA C-MAPSS Dataset  ·  Random Forest Models  ·  Real-time RUL & Failure Type Prediction",
             Inches(0.5), Inches(4.0), Inches(12), Inches(0.6),
             font_size=16, color=RGBColor(0xa8, 0xc8, 0xf0))

add_text_box(sl, "Machine Learning · Predictive Analytics · Streamlit Dashboard",
             Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
             font_size=13, color=RGBColor(0x6b, 0x9a, 0xc8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Problem Statement", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

points = [
    "🔧  Unplanned engine failures cause costly downtime and safety risks in aviation.",
    "📉  Traditional time-based maintenance leads to over-servicing or under-servicing.",
    "🎯  Goal: Predict Remaining Useful Life (RUL) and failure type before failure occurs.",
    "📊  Dataset: NASA C-MAPSS — 104 turbofan engine run-to-failure simulations.",
    "🤖  Solution: Train ML models on sensor history to enable condition-based maintenance.",
]
bullet_box(sl, points, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
           font_size=18, color=WHITE, spacing=Pt(14))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Dataset Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Dataset Overview", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

# stat cards row 1
cards = [
    ("Source",          "NASA C-MAPSS"),
    ("Total Engines",   "104 rows"),
    ("Unique Engines",  "26 × 4 types"),
    ("Sensors",         "21 sensors"),
]
for i, (t, v) in enumerate(cards):
    add_card(sl, Inches(0.5 + i * 3.1), Inches(1.3), Inches(2.9), Inches(1.1), t, v)

# stat cards row 2
cards2 = [
    ("Op. Settings",    "3 settings"),
    ("Failure Types",   "4 classes"),
    ("RUL Range",       "0 – 347 cycles"),
    ("Features",        "144 engineered"),
]
for i, (t, v) in enumerate(cards2):
    add_card(sl, Inches(0.5 + i * 3.1), Inches(2.6), Inches(2.9), Inches(1.1), t, v)

add_text_box(sl, "Failure Types", Inches(0.5), Inches(4.0), Inches(4), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_BLU)
failure_types = [
    "Failure_1 — High Bypass Fan degradation",
    "Failure_2 — Low Pressure Compressor fault",
    "Failure_3 — High Pressure Compressor degradation",
    "Failure_4 — Turbine degradation",
]
bullet_box(sl, failure_types, Inches(0.6), Inches(4.45), Inches(6), Inches(2.5),
           font_size=15, color=WHITE, spacing=Pt(8))

add_text_box(sl, "Feature Engineering", Inches(6.8), Inches(4.0), Inches(4), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_BLU)
fe_points = [
    "Rolling avg over 30, 60, 90 cycles",
    "Min, Max, Variance per signal",
    "24 signals × 6 stats = 144 features",
]
bullet_box(sl, fe_points, Inches(6.9), Inches(4.45), Inches(5.8), Inches(2.0),
           font_size=15, color=WHITE, spacing=Pt(8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Methodology / Pipeline
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "ML Pipeline", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

steps = [
    ("1", "Raw Sensor Data",       "engine_id, cycle, 21 sensors, 3 op settings"),
    ("2", "Feature Engineering",   "144 rolling statistics per engine"),
    ("3", "Imputation",            "SimpleImputer (mean strategy) for missing values"),
    ("4", "RUL Regression",        "Random Forest Regressor → predicts cycles remaining"),
    ("5", "Failure Classification","Random Forest Classifier → predicts failure type"),
    ("6", "Dashboard Output",      "Gauge, alerts, fleet charts, CSV export"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.25 + i * 0.88)
    add_rect(sl, Inches(0.5), y, Inches(0.55), Inches(0.65), GREEN)
    add_text_box(sl, num, Inches(0.5), y + Inches(0.08), Inches(0.55), Inches(0.5),
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.15), y, Inches(11.5), Inches(0.65), BLUE)
    add_text_box(sl, title, Inches(1.25), y + Inches(0.04), Inches(3.2), Inches(0.35),
                 font_size=15, bold=True, color=WHITE)
    add_text_box(sl, desc, Inches(4.5), y + Inches(0.04), Inches(8.0), Inches(0.55),
                 font_size=13, color=LIGHT_BLU)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Model Performance
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Model Performance", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

# ── Classification table ──────────────────────────────────────────────────────
add_text_box(sl, "Classification — Failure Type", Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_BLU)

cls_headers = ["Model", "Accuracy", "Precision", "Recall", "F1", "CV Mean"]
cls_rows = [
    ["Random Forest",     "80.95%", "81.20%", "80.95%", "80.85%", "80.62%"],
    ["Logistic Regression","80.95%","80.10%", "80.95%", "80.30%", "83.67%"],
]
col_w = [Inches(2.0), Inches(0.9), Inches(0.9), Inches(0.8), Inches(0.7), Inches(0.9)]
x0, y0 = Inches(0.5), Inches(1.65)
row_h = Inches(0.42)

# header
x = x0
for j, h in enumerate(cls_headers):
    add_rect(sl, x, y0, col_w[j], row_h, NAVY)
    add_text_box(sl, h, x + Inches(0.05), y0 + Inches(0.06), col_w[j], row_h,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w[j]

for i, row in enumerate(cls_rows):
    y = y0 + row_h * (i + 1)
    fill = GREEN if i == 0 else LIGHT_BLU
    txt_c = WHITE if i == 0 else NAVY
    x = x0
    for j, cell in enumerate(row):
        add_rect(sl, x, y, col_w[j], row_h, fill)
        add_text_box(sl, cell, x + Inches(0.05), y + Inches(0.06), col_w[j], row_h,
                     font_size=12, bold=(j == 0), color=txt_c, align=PP_ALIGN.CENTER)
        x += col_w[j]

# ── Regression table ──────────────────────────────────────────────────────────
add_text_box(sl, "Regression — RUL Prediction", Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_BLU)

reg_headers = ["Model", "MAE", "RMSE", "R²", "vs Baseline"]
reg_rows = [
    ["Random Forest",    "34.87", "40.87", "0.83", "22.8% better"],
    ["Linear Regression","45.17", "60.42", "0.71", "—"],
]
col_w2 = [Inches(2.0), Inches(0.8), Inches(0.85), Inches(0.7), Inches(1.5)]
x0b = Inches(7.0)

x = x0b
for j, h in enumerate(reg_headers):
    add_rect(sl, x, y0, col_w2[j], row_h, NAVY)
    add_text_box(sl, h, x + Inches(0.05), y0 + Inches(0.06), col_w2[j], row_h,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w2[j]

for i, row in enumerate(reg_rows):
    y = y0 + row_h * (i + 1)
    fill = GREEN if i == 0 else LIGHT_BLU
    txt_c = WHITE if i == 0 else NAVY
    x = x0b
    for j, cell in enumerate(row):
        add_rect(sl, x, y, col_w2[j], row_h, fill)
        add_text_box(sl, cell, x + Inches(0.05), y + Inches(0.06), col_w2[j], row_h,
                     font_size=12, bold=(j == 0), color=txt_c, align=PP_ALIGN.CENTER)
        x += col_w2[j]

# key takeaways
divider(sl, Inches(3.6), GREEN)
add_text_box(sl, "Key Takeaways", Inches(0.5), Inches(3.7), Inches(5), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_BLU)
takeaways = [
    "✅  Random Forest outperforms baseline on both tasks",
    "✅  Regressor achieves MAE of 34.87 cycles (22.8% improvement)",
    "✅  Classifier reaches 80.95% accuracy across 4 failure types",
    "✅  High correlation (0.95) between true and predicted RUL",
]
bullet_box(sl, takeaways, Inches(0.6), Inches(4.15), Inches(12), Inches(2.8),
           font_size=15, color=WHITE, spacing=Pt(10))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Top Features
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Top 10 Features — RUL Prediction", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

features = [
    ("sensor_15_avg_90", 0.213),
    ("sensor_15_min",    0.136),
    ("sensor_15_avg_60", 0.083),
    ("sensor_15_avg_30", 0.077),
    ("sensor_21_max",    0.041),
    ("sensor_4_min",     0.034),
    ("sensor_12_max",    0.026),
    ("sensor_3_min",     0.023),
    ("sensor_17_avg_60", 0.020),
    ("sensor_7_max",     0.016),
]

bar_max = Inches(7.5)
for i, (name, imp) in enumerate(features):
    y = Inches(1.25 + i * 0.57)
    add_text_box(sl, name, Inches(0.5), y, Inches(2.2), Inches(0.45),
                 font_size=13, color=LIGHT_BLU)
    bar_w = bar_max * (imp / 0.213)
    bar_color = GREEN if i == 0 else (BLUE if i < 4 else RGBColor(0x3b, 0x82, 0xf6))
    add_rect(sl, Inches(2.8), y + Inches(0.06), bar_w, Inches(0.33), bar_color)
    add_text_box(sl, f"{imp:.3f}", Inches(2.8) + bar_w + Inches(0.1), y + Inches(0.06),
                 Inches(0.8), Inches(0.33), font_size=12, color=WHITE)

add_text_box(sl, "💡  Sensor 15 features alone account for 51.3% of total predictive power",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             font_size=14, color=RGBColor(0xa8, 0xc8, 0xf0))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Dashboard Demo
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Streamlit Dashboard", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

features_dash = [
    ("📂  Upload & Predict",    "Upload raw sensor CSV or pre-engineered feature CSV.\nAuto-detects format and runs predictions instantly."),
    ("🔴🟠🟢  Status Alerts",   "CRITICAL (<50 cycles), WARNING (<150 cycles),\nHEALTHY (≥150 cycles) with colour-coded alerts."),
    ("📊  Fleet View",          "Bar chart of RUL per engine, pie chart of failure\ntype distribution, colour-coded results table."),
    ("⬇️  Export",              "Download full prediction results as CSV for\nreporting and further analysis."),
    ("📈  Model Info Tab",      "Performance tables, feature importance chart,\ndataset statistics and pipeline explanation."),
]

for i, (title, desc) in enumerate(features_dash):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.3 + row * 1.9)
    add_rect(sl, x, y, Inches(6.0), Inches(1.65), BLUE)
    add_text_box(sl, title, x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.45),
                 font_size=15, bold=True, color=WHITE)
    add_text_box(sl, desc, x + Inches(0.2), y + Inches(0.55), Inches(5.6), Inches(1.0),
                 font_size=13, color=LIGHT_BLU)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – RUL Status Zones
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "RUL Status Zones", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

zones = [
    (RED,    "🔴  CRITICAL",  "RUL < 50 cycles",   "Immediate maintenance required.\nEngine at high risk of failure."),
    (ORANGE, "🟠  WARNING",   "50 ≤ RUL < 150",    "Schedule maintenance soon.\nMonitor closely."),
    (GREEN,  "🟢  HEALTHY",   "RUL ≥ 150 cycles",  "Engine operating normally.\nNo immediate action needed."),
]

for i, (color, label, rul_range, desc) in enumerate(zones):
    x = Inches(0.5 + i * 4.2)
    add_rect(sl, x, Inches(1.3), Inches(3.9), Inches(4.5), color)
    add_text_box(sl, label, x + Inches(0.2), Inches(1.5), Inches(3.5), Inches(0.7),
                 font_size=22, bold=True, color=WHITE)
    add_text_box(sl, rul_range, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.5),
                 font_size=16, color=WHITE)
    add_text_box(sl, desc, x + Inches(0.2), Inches(2.9), Inches(3.5), Inches(1.5),
                 font_size=14, color=WHITE)

add_text_box(sl, "Demo CSV covers all 3 zones: 1 CRITICAL · 2 WARNING · 4 HEALTHY engines",
             Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
             font_size=14, color=RGBColor(0xa8, 0xc8, 0xf0), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Conclusion
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, GREEN)

add_text_box(sl, "Conclusion & Future Work", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
divider(sl, Inches(1.1), GREEN)

add_text_box(sl, "Achievements", Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=18, bold=True, color=LIGHT_BLU)
achievements = [
    "✅  End-to-end predictive maintenance pipeline",
    "✅  Random Forest Regressor — MAE 34.87 cycles",
    "✅  Random Forest Classifier — 80.95% accuracy",
    "✅  144 engineered features from raw sensor data",
    "✅  Interactive Streamlit dashboard with fleet view",
]
bullet_box(sl, achievements, Inches(0.6), Inches(1.65), Inches(5.8), Inches(3.0),
           font_size=15, color=WHITE, spacing=Pt(10))

add_text_box(sl, "Future Work", Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
             font_size=18, bold=True, color=LIGHT_BLU)
future = [
    "🔮  Deep learning models (LSTM, Transformer)",
    "🔮  Real-time streaming sensor data ingestion",
    "🔮  Anomaly detection layer",
    "🔮  Multi-fleet comparative analytics",
    "🔮  Maintenance scheduling integration",
]
bullet_box(sl, future, Inches(7.1), Inches(1.65), Inches(5.8), Inches(3.0),
           font_size=15, color=WHITE, spacing=Pt(10))

divider(sl, Inches(5.1), GREEN)
add_text_box(sl, "\"Predict before it breaks — save cost, save lives.\"",
             Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.7),
             font_size=20, bold=True, color=LIGHT_BLU, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "Turbofan_Predictive_Maintenance.pptx"
prs.save(out)
print(f"Saved: {out}")
